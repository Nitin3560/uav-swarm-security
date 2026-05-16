from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import yaml
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/nitin/Desktop/failure/paper_sim")
RESULTS_ROOT = ROOT / "results_final_hybrid_10seed"
CONFIG_PATH = ROOT / "configs" / "final_hybrid.yaml"
OUTPUT_DIR = ROOT / "presentation"
ASSET_DIR = OUTPUT_DIR / "assets"
PPTX_PATH = OUTPUT_DIR / "uav_swarm_fault_tolerance_technical_review.pptx"

SCENARIOS = ["nominal", "wind", "sensor", "comm"]
CONTROLLERS = ["pid", "generic", "failure_aware"]
COLORS = {
    "pid": "#6c757d",
    "generic": "#f4a261",
    "failure_aware": "#1d4ed8",
}


@dataclass
class SlideTheme:
    navy: RGBColor = RGBColor(18, 34, 66)
    blue: RGBColor = RGBColor(29, 78, 216)
    teal: RGBColor = RGBColor(18, 130, 162)
    orange: RGBColor = RGBColor(230, 126, 34)
    red: RGBColor = RGBColor(184, 51, 51)
    green: RGBColor = RGBColor(26, 127, 55)
    gray: RGBColor = RGBColor(89, 99, 110)
    light_bg: RGBColor = RGBColor(245, 247, 250)


THEME = SlideTheme()


def load_config() -> dict:
    with CONFIG_PATH.open("r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def load_aggregate(scenario: str) -> pd.DataFrame:
    return pd.read_csv(RESULTS_ROOT / scenario / "aggregate_summary.csv")


def load_trace(scenario: str, controller: str, seed: int = 1) -> pd.DataFrame:
    return pd.read_csv(RESULTS_ROOT / scenario / "csv" / f"{controller}_seed{seed}.csv")


def ensure_dirs() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def style_axes(ax: plt.Axes) -> None:
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(True, alpha=0.22, linewidth=0.8)
    ax.tick_params(labelsize=9)


def save_architecture_diagram() -> Path:
    path = ASSET_DIR / "architecture_diagram.png"
    fig, ax = plt.subplots(figsize=(12, 6))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 7)
    ax.axis("off")

    def box(x: float, y: float, w: float, h: float, title: str, body: str, fc: str) -> None:
        patch = FancyBboxPatch(
            (x, y),
            w,
            h,
            boxstyle="round,pad=0.03,rounding_size=0.15",
            linewidth=2,
            edgecolor="#1f2937",
            facecolor=fc,
        )
        ax.add_patch(patch)
        ax.text(x + w / 2, y + h - 0.35, title, ha="center", va="top", fontsize=15, weight="bold")
        ax.text(x + 0.2, y + h - 0.8, body, ha="left", va="top", fontsize=11)

    box(
        0.5,
        3.9,
        2.6,
        2.2,
        "Reference Generator",
        "Circular reference p_ref(t), v_ref(t)\nSquare formation offsets\nMission schedule: settle / nominal / fault / recovery",
        "#dbeafe",
    )
    box(
        3.7,
        3.9,
        2.6,
        2.2,
        "Failure-Aware Supervisor",
        "Diagnosis-driven gain scheduling\nBounded speed / formation / consensus scaling\nConnectivity bias and wind compensation",
        "#dcfce7",
    )
    box(
        3.7,
        0.9,
        2.6,
        2.2,
        "Fault Diagnosis + Estimation",
        "Residual confidences: sensor / wind / comm\nConstant-velocity prediction\nInnovation gating and filtered state",
        "#fef3c7",
    )
    box(
        6.9,
        3.9,
        2.6,
        2.2,
        "PID Inner Loop",
        "u_i = Kp e_i + Ki integral(e_i) + Kd dot(e_i)\nAnti-windup saturation handling\nPer-UAV acceleration commands",
        "#fce7f3",
    )
    box(
        10.1,
        3.9,
        1.4,
        2.2,
        "Swarm Dynamics",
        "4 UAVs\nDrag + saturation\nWind enters as exogenous acceleration",
        "#e0f2fe",
    )
    box(
        10.1,
        0.9,
        1.4,
        2.2,
        "Fault Layer",
        "Wind gusts\nSensor bias / drift / spikes\nPacket drop / delay / isolation",
        "#fee2e2",
    )

    arrows = [
        ((3.1, 5.0), (3.7, 5.0)),
        ((5.0, 3.9), (5.0, 3.1)),
        ((6.3, 5.0), (6.9, 5.0)),
        ((9.5, 5.0), (10.1, 5.0)),
        ((10.8, 3.9), (10.8, 3.1)),
        ((10.1, 2.0), (6.3, 2.0)),
        ((6.3, 2.0), (6.3, 4.45)),
    ]
    for a, b in arrows:
        ax.add_patch(
            FancyArrowPatch(
                a,
                b,
                arrowstyle="-|>",
                mutation_scale=18,
                linewidth=2,
                color="#334155",
            )
        )

    ax.text(7.9, 2.45, "measured / degraded observations", fontsize=10, color="#334155")
    ax.text(4.2, 3.35, "diagnosis state", fontsize=10, color="#334155")
    ax.text(7.35, 5.4, "bounded supervisory modification", fontsize=10, color="#334155")
    fig.tight_layout()
    fig.savefig(path, dpi=220, bbox_inches="tight")
    plt.close(fig)
    return path


def save_tracking_grid() -> Path:
    path = ASSET_DIR / "tracking_error_grid.png"
    fig, axes = plt.subplots(2, 2, figsize=(13.5, 8), sharex=True)
    axes = axes.ravel()
    fault_start, fault_end = 20.0, 30.0
    for ax, scenario in zip(axes, SCENARIOS):
        for controller in CONTROLLERS:
            df = load_trace(scenario, controller, seed=1)
            ax.plot(df["t"], df["mean_err_nominal_m"], linewidth=2.0, color=COLORS[controller], label=controller.replace("_", " ").title())
        ax.axvspan(fault_start, fault_end, color="#fde68a", alpha=0.35)
        ax.axvline(10.0, color="#94a3b8", linestyle="--", linewidth=1)
        ax.axvline(20.0, color="#94a3b8", linestyle="--", linewidth=1)
        ax.axvline(30.0, color="#94a3b8", linestyle="--", linewidth=1)
        ax.set_title(scenario.title(), fontsize=13, weight="bold")
        ax.set_ylabel("Mean Tracking Error (m)")
        ax.set_xlabel("Time (s)")
        style_axes(ax)
    handles, labels = axes[0].get_legend_handles_labels()
    fig.legend(handles, labels, loc="upper center", ncol=3, frameon=False, fontsize=10)
    fig.suptitle("Tracking Error Timeline from Logged Seed-1 Traces", fontsize=16, weight="bold", y=0.98)
    fig.tight_layout(rect=[0, 0, 1, 0.95])
    fig.savefig(path, dpi=220, bbox_inches="tight")
    plt.close(fig)
    return path


def save_baseline_comparison() -> Path:
    path = ASSET_DIR / "baseline_comparison.png"
    metrics = ["fault_error_m", "rmse_m", "recovery_time_s"]
    metric_titles = ["Fault-Window Mean Error (m)", "RMSE (m)", "Recovery Time (s)"]
    fig, axes = plt.subplots(1, 3, figsize=(14, 4.5))
    x = np.arange(len(SCENARIOS))
    width = 0.24
    for ax, metric, metric_title in zip(axes, metrics, metric_titles):
        for idx, controller in enumerate(CONTROLLERS):
            values = []
            for scenario in SCENARIOS:
                agg = load_aggregate(scenario).set_index("controller_label")
                values.append(float(agg.loc[controller, metric]))
            ax.bar(x + (idx - 1) * width, values, width=width, color=COLORS[controller], label=controller.replace("_", " ").title())
        ax.set_xticks(x)
        ax.set_xticklabels([s.title() for s in SCENARIOS], rotation=0)
        ax.set_title(metric_title, fontsize=12, weight="bold")
        style_axes(ax)
    axes[0].legend(frameon=False, fontsize=9)
    fig.suptitle("Aggregate Comparison Across Controllers", fontsize=16, weight="bold")
    fig.tight_layout(rect=[0, 0, 1, 0.94])
    fig.savefig(path, dpi=220, bbox_inches="tight")
    plt.close(fig)
    return path


def save_diagnosis_plot() -> Path:
    path = ASSET_DIR / "diagnosis_response.png"
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.5))

    wind_df = load_trace("wind", "failure_aware", seed=1)
    axes[0].plot(wind_df["t"], wind_df["wind_confidence"], color="#1d4ed8", linewidth=2, label="Wind confidence")
    axes[0].plot(wind_df["t"], wind_df["disturbance_estimate_norm"], color="#ef4444", linewidth=1.8, label="Disturbance estimate norm")
    axes[0].axvspan(20.0, 30.0, color="#fde68a", alpha=0.35)
    axes[0].set_title("Wind Scenario: Diagnosis + Observer", fontsize=12, weight="bold")
    axes[0].set_xlabel("Time (s)")
    axes[0].set_ylabel("Signal magnitude")
    style_axes(axes[0])
    axes[0].legend(frameon=False, fontsize=9)

    sensor_df = load_trace("sensor", "failure_aware", seed=1)
    axes[1].plot(sensor_df["t"], sensor_df["sensor_confidence"], color="#1d4ed8", linewidth=2, label="Sensor confidence")
    axes[1].plot(sensor_df["t"], sensor_df["innovation_norm_m"], color="#f97316", linewidth=1.8, label="Innovation norm")
    axes[1].plot(sensor_df["t"], sensor_df["measurement_accept_rate"], color="#16a34a", linewidth=1.8, label="Accept rate")
    axes[1].axvspan(20.0, 30.0, color="#fde68a", alpha=0.35)
    axes[1].set_title("Sensor Scenario: Estimation Signals", fontsize=12, weight="bold")
    axes[1].set_xlabel("Time (s)")
    axes[1].set_ylabel("Signal magnitude")
    style_axes(axes[1])
    axes[1].legend(frameon=False, fontsize=9)

    fig.suptitle("Internal Diagnosis / Estimation Signals from Logged Runs", fontsize=15, weight="bold")
    fig.tight_layout(rect=[0, 0, 1, 0.93])
    fig.savefig(path, dpi=220, bbox_inches="tight")
    plt.close(fig)
    return path


def add_slide_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.0), Inches(0.55))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = THEME.navy
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.78), Inches(11.8), Inches(0.35))
        sub_tf = sub_box.text_frame
        sub_p = sub_tf.paragraphs[0]
        sub_run = sub_p.add_run()
        sub_run.text = subtitle
        sub_run.font.size = Pt(13)
        sub_run.font.color.rgb = THEME.gray


def add_bullets(slide, left: float, top: float, width: float, height: float, bullets: Iterable[str], font_size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = THEME.navy
        p.space_after = Pt(8)


def add_text_block(slide, left: float, top: float, width: float, height: float, text: str, font_size: int = 18, color: RGBColor | None = None, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color or THEME.navy
    run.font.bold = bold


def add_equation_box(slide, left: float, top: float, width: float, height: float, lines: list[str]) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(246, 248, 252)
    shape.line.color.rgb = RGBColor(203, 213, 225)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Courier New"
        p.font.size = Pt(19 if idx == 0 else 17)
        p.font.color.rgb = THEME.navy


def add_table(slide, left: float, top: float, width: float, height: float, df: pd.DataFrame, font_size: int = 11) -> None:
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    for c, col in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME.navy
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(font_size)
                run.font.color.rgb = RGBColor(255, 255, 255)
    for r in range(df.shape[0]):
        for c in range(cols):
            cell = table.cell(r + 1, c)
            val = df.iat[r, c]
            cell.text = str(val)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 250, 252)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = THEME.navy


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(6.92), Inches(12.0), Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = THEME.gray


def build_result_table() -> pd.DataFrame:
    rows = []
    for scenario in SCENARIOS:
        agg = load_aggregate(scenario)
        for _, row in agg.iterrows():
            rows.append(
                {
                    "Scenario": scenario.title(),
                    "Controller": row["controller_label"].replace("_", " ").title(),
                    "Fault Err (m)": f"{row['fault_error_m']:.3f}",
                    "RMSE (m)": f"{row['rmse_m']:.3f}",
                    "Recovery (s)": "N/A" if pd.isna(row["recovery_time_s"]) else f"{row['recovery_time_s']:.2f}",
                    "CI95": f"{row['fault_error_ci95']:.4f}",
                }
            )
    return pd.DataFrame(rows)


def build_comparison_table() -> pd.DataFrame:
    rows = []
    for scenario in ["wind", "sensor", "comm"]:
        agg = load_aggregate(scenario).set_index("controller_label")
        pid_err = float(agg.loc["pid", "fault_error_m"])
        generic_err = float(agg.loc["generic", "fault_error_m"])
        hybrid_err = float(agg.loc["failure_aware", "fault_error_m"])
        rows.append(
            {
                "Scenario": scenario.title(),
                "Hybrid vs PID": f"{(pid_err - hybrid_err) / pid_err * 100:.1f}%",
                "Hybrid vs Generic": f"{(generic_err - hybrid_err) / generic_err * 100:.1f}%",
                "Hybrid Recovery (s)": "N/A" if pd.isna(agg.loc["failure_aware", "recovery_time_s"]) else f"{agg.loc['failure_aware', 'recovery_time_s']:.2f}",
            }
        )
    return pd.DataFrame(rows)


def build_presentation() -> Path:
    ensure_dirs()
    cfg = load_config()
    architecture_path = save_architecture_diagram()
    tracking_path = save_tracking_grid()
    comparison_path = save_baseline_comparison()
    diagnosis_path = save_diagnosis_plot()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME.light_bg
    bg.line.fill.background()
    add_text_block(slide, 0.7, 1.0, 11.8, 0.8, "Diagnosis-Driven Hybrid UAV Swarm Control", 28, THEME.navy, True)
    add_text_block(slide, 0.72, 1.75, 8.5, 0.4, "Fault-Tolerant UAV Swarm Control", 20, THEME.blue, True)
    add_text_block(slide, 0.72, 2.35, 9.5, 1.0, "Technical presentation of the standalone paper_sim simulator,\ncontroller stack, fault models, metrics, and final hybrid results.", 18, THEME.gray)
    add_text_block(slide, 0.72, 6.3, 3.2, 0.3, "Author: Nitin", 16, THEME.gray)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Motivation", "Why fault-tolerant swarm control matters")
    add_bullets(
        slide,
        0.65,
        1.25,
        6.0,
        4.8,
        [
            "Real UAV swarms face coupled failures: wind disturbance, corrupted onboard sensing, and network degradation.",
            "A classical PID loop is purely local and reactive: it does not know whether the dominant issue is sensing, environment, or communication.",
            "Communication loss is especially dangerous because low-level control may stay stable while formation-level coordination silently degrades.",
            "The study goal is not only trajectory following, but resilient tracking, formation preservation, and recovery under faults.",
        ],
    )
    add_equation_box(
        slide,
        7.0,
        1.55,
        5.5,
        2.1,
        [
            "Observed state = true state + disturbance + fault corruption",
            "x_meas = x_true + eta_noise + b_sensor + dropout/staleness",
        ],
    )
    add_text_block(slide, 7.05, 4.0, 5.4, 1.8, "Classical PID is insufficient because it acts on local error only.\nIt cannot infer whether an error should be corrected aggressively, filtered, or coordinated across neighbors.", 17, THEME.navy)
    add_footer(slide, "Standalone simulator: /Users/nitin/Desktop/failure/paper_sim")

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Problem Statement", "Formal objective and operating constraints")
    add_equation_box(
        slide,
        0.65,
        1.3,
        5.5,
        2.2,
        [
            "e_i(t) = || p_i(t) - p_i^ref(t) ||_2",
            "minimize  J = w1 * mean_t mean_i e_i(t) + w2 * formation_error + w3 * recovery_time",
        ],
    )
    add_bullets(
        slide,
        0.75,
        3.9,
        5.7,
        2.5,
        [
            "Maintain 4-UAV square formation while following a circular mission trajectory.",
            "Remain stable under active faults during 20-30 s and recover during 30-40 s.",
            "Bound control inputs under acceleration and speed saturation constraints.",
        ],
        font_size=17,
    )
    add_text_block(slide, 6.9, 1.45, 5.2, 0.4, "Primary research question", 18, THEME.blue, True)
    add_text_block(slide, 6.9, 1.95, 5.6, 1.2, "Can a diagnosis-driven supervisory layer, placed above a PID inner loop, improve robustness to wind, sensor, and communication faults without destabilizing nominal operation?", 18, THEME.navy)
    add_bullets(
        slide,
        6.95,
        3.65,
        5.6,
        2.5,
        [
            "Tracking objective: minimize mean XY trajectory error.",
            "Coordination objective: preserve spacing/connectivity.",
            "Resilience objective: reduce error spike and recovery time across seeds.",
        ],
        font_size=17,
    )
    add_footer(slide, "Optimization is evaluated empirically through multi-seed simulation rather than solved as an offline optimal-control program.")

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "System Overview", "Layered architecture used in the final hybrid stack")
    slide.shapes.add_picture(str(architecture_path), Inches(0.65), Inches(1.15), width=Inches(12.0))
    add_footer(slide, "Supervisor modifies bounded coordination parameters; PID remains the stabilizing low-level controller.")

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Technical Analysis: Control Layer", "PID inner loop and its limitations")
    add_equation_box(
        slide,
        0.7,
        1.35,
        5.4,
        1.8,
        [
            "u_i(t) = Kp e_i(t) + Ki integral e_i(t) dt + Kd dot(e_i)(t)",
            "e_i(t) = p_i^cmd(t) - p_i^meas(t),   dot(e_i) = v_i^cmd(t) - v_i^meas(t)",
        ],
    )
    add_bullets(
        slide,
        0.72,
        3.45,
        5.5,
        2.8,
        [
            "Implemented per UAV in controllers.py with separate XY and Z gains, saturation limits, and anti-windup correction.",
            "Stability intuition: proportional action pulls the vehicle toward the command state, derivative damping resists overshoot, and integral action removes steady-state bias.",
            "The PID loop is effective when measurements are trustworthy and disturbances are modest.",
        ],
        font_size=16,
    )
    add_bullets(
        slide,
        6.65,
        1.55,
        5.6,
        4.8,
        [
            "Limitation 1: if sensing is corrupted, the controller acts on wrong state information.",
            "Limitation 2: if packets are dropped or stale, PID has no explicit network awareness.",
            "Limitation 3: under large wind, PID rejects disturbance only indirectly through tracking error.",
            "Therefore PID is retained as a stabilizing baseline, while higher-level supervision handles context and fault mode.",
        ],
        font_size=17,
    )
    add_footer(slide, "Code reference: PositionPID in /Users/nitin/Desktop/failure/paper_sim/controllers.py")

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Technical Analysis: Supervisor", "Diagnosis-driven bounded reconfiguration")
    add_equation_box(
        slide,
        0.65,
        1.4,
        5.7,
        1.7,
        [
            "p_ref,new = p_ref + Delta_ref,   ||Delta_ref|| <= Delta_max",
            "theta_sup = { speed_scale, formation_scale, consensus_scale, connectivity_bias }",
        ],
    )
    add_bullets(
        slide,
        0.7,
        3.45,
        5.8,
        2.9,
        [
            "The supervisor never replaces PID; it modifies the commanded trajectory and coupling terms through bounded scalars and biases.",
            "Activation input is the diagnosis snapshot: active fault label plus confidence scores in [0, 1].",
            "Recovery uses a blend-back schedule to avoid abrupt gain switching after the fault window clears.",
        ],
        font_size=16,
    )
    add_text_block(slide, 6.7, 1.45, 5.4, 0.35, "Mode-specific behavior", 18, THEME.blue, True)
    add_bullets(
        slide,
        6.72,
        1.95,
        5.6,
        3.8,
        [
            "Wind: reduce commanded speed, increase damping, and add bias against measured group drift.",
            "Sensor: rely more on filtered state than raw measurement; soften consensus sensitivity to corrupted observations.",
            "Communication: reduce dependence on degraded neighbors and inject connectivity-preserving bias toward the swarm center.",
            "Nominal recovery: exponentially blend supervisory parameters back to 1.0-equivalent values.",
        ],
        font_size=16,
    )
    add_footer(slide, "Code reference: FailureAwareSupervisor.step(...) in /Users/nitin/Desktop/failure/paper_sim/controllers.py")

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Fault Modeling", "How the simulator injects disturbances and failures")
    add_equation_box(
        slide,
        0.7,
        1.3,
        4.0,
        2.7,
        [
            "Wind: a_wind(t) = a_base + a_gust(t) + a_var(t) + epsilon",
            "Sensor: p_meas = p_true + b + d(t) + n + spikes / freeze / dropout",
            "Comm: x_ij^recv(t) = x_j(t - tau) or None with packet loss",
        ],
    )
    fault_table = pd.DataFrame(
        [
            ["Wind", "gust + constant + varying", "base [0.55,0,0], gust [0.85,0,0], interval 1.5 s"],
            ["Sensor", "mixed corruption", "bias [1.0,-0.5], drift 0.04 m/s, spikes 0.60 m"],
            ["Comm", "degraded network", "85% dropout, 12-step delay, one-agent isolation"],
        ],
        columns=["Fault", "Model", "Scenario parameters"],
    )
    add_table(slide, 5.0, 1.45, 7.3, 2.2, fault_table, font_size=10)
    add_bullets(
        slide,
        0.75,
        4.35,
        11.5,
        1.8,
        [
            "All fault processes are seed-controlled in faults.py, so the 10-seed comparisons are reproducible.",
            "Wind acts as exogenous acceleration in the dynamics layer; sensor and communication faults corrupt the information path rather than the physical state itself.",
        ],
        font_size=17,
    )
    add_footer(slide, "Code reference: FaultModel in /Users/nitin/Desktop/failure/paper_sim/faults.py")

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Experimental Setup", "Simulation protocol and data collection")
    add_bullets(
        slide,
        0.7,
        1.35,
        5.8,
        4.9,
        [
            f"Swarm size: {cfg['sim']['num_drones']} UAVs in a square formation with {cfg['formation']['spacing_m']:.1f} m spacing.",
            f"Reference: circular trajectory, radius {cfg['trajectory']['radius_m']:.1f} m, period {cfg['trajectory']['period_s']:.1f} s, altitude {cfg['trajectory']['altitude_m']:.1f} m.",
            f"Simulation frequency: {cfg['sim']['freq_hz']} Hz, control frequency: {cfg['sim']['ctrl_hz']} Hz.",
            "Compared controllers: PID, generic supervisor, diagnosis-driven hybrid.",
            "Main evaluation: 10 seeded runs per scenario under identical schedule and parameterization.",
        ],
        font_size=17,
    )
    timeline = pd.DataFrame(
        [
            ["0-10 s", "Settling", "Swarm converges to reference orbit and formation"],
            ["10-20 s", "Nominal", "Reference tracking before injected fault"],
            ["20-30 s", "Fault active", "Wind, sensor, or communication degradation enabled"],
            ["30-40 s", "Recovery", "Fault removed, controller must return toward nominal tracking"],
        ],
        columns=["Window", "Phase", "Purpose"],
    )
    add_table(slide, 6.8, 1.45, 5.4, 2.5, timeline, font_size=11)
    add_text_block(slide, 6.9, 4.35, 5.2, 1.2, "Validity checks\n• Same environment and mission for all controllers\n• Same seeds across baselines\n• Metrics computed from logged control-step traces", 17, THEME.navy)
    add_footer(slide, "Experiment runner: /Users/nitin/Desktop/failure/paper_sim/run_study.py and run_matrix.py")

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Metrics", "How performance is computed from logged traces")
    add_equation_box(
        slide,
        0.7,
        1.25,
        5.7,
        2.6,
        [
            "mean error(t_k) = (1/N) * sum_i || p_i(t_k) - p_i^ref(t_k) ||_2",
            "RMSE = sqrt( (1/T) * sum_k mean_error(t_k)^2 )",
            "degradation % = (fault_mean - pre_fault_mean) / pre_fault_mean * 100",
        ],
    )
    add_equation_box(
        slide,
        6.75,
        1.25,
        5.3,
        2.6,
        [
            "recovery_time = first t >= t_fault s.t. e(t:t+tau) <= gamma * e_pre",
            "max formation deformation = max_t max_{i<j} | d_ij(t) - d_ij^des |",
            "stable run = 1 if max_t e(t) < 2.5 * safety_threshold",
        ],
    )
    add_bullets(
        slide,
        0.8,
        4.15,
        11.4,
        2.0,
        [
            "pre-, fault-, and post-fault means are computed over [10,20), [20,30), and [30,40) seconds respectively.",
            "Recovery and settling times are threshold-based with a 2 s sustain requirement, using the configuration values in final_hybrid.yaml.",
            "run_matrix.py aggregates mean, standard deviation, and 95% confidence interval across 10 seeds.",
        ],
        font_size=17,
    )
    add_footer(slide, "Metric definitions come directly from /Users/nitin/Desktop/failure/paper_sim/metrics.py and run_study.py")

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Results Tables", "Final 10-seed aggregate results from the unified hybrid package")
    result_df = build_result_table()
    add_table(slide, 0.55, 1.15, 12.2, 5.6, result_df, font_size=10)
    add_footer(slide, "Source: /Users/nitin/Desktop/failure/paper_sim/results_final_hybrid_10seed/*/aggregate_summary.csv")

    # Slide 11
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Results Plots", "Tracking error from real logged traces with fault window highlighted")
    slide.shapes.add_picture(str(tracking_path), Inches(0.45), Inches(1.05), width=Inches(12.4))
    add_footer(slide, "Shaded region: active fault window [20, 30) s. Data source: seed-1 CSV logs for each controller/scenario.")

    # Slide 12
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Technical Interpretation", "Why the hybrid behaves the way it does")
    slide.shapes.add_picture(str(diagnosis_path), Inches(6.7), Inches(1.25), width=Inches(5.7))
    add_bullets(
        slide,
        0.7,
        1.45,
        5.7,
        4.9,
        [
            "Communication is the clearest strength because the supervisory layer has information that pure PID does not: packet quality, stale neighbor state, and connectivity biasing.",
            "Sensor faults remain hard because corrupted state hurts both the low-level controller and the diagnosis layer; filtered-state handling helps, but estimation is still the limiting factor.",
            "Wind becomes a hybrid advantage in the final package because damping and drift compensation act before large tracking error accumulates.",
            "Nominal performance also improves in the final hybrid package, so the final story is not only fault tolerance but overall better command shaping.",
        ],
        font_size=16,
    )
    add_footer(slide, "Interpretation is supported by the logged diagnosis confidences, innovation traces, and disturbance-estimate channels.")

    # Slide 13
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Comparison with Baselines", "Aggregate view: PID vs generic supervisor vs hybrid")
    slide.shapes.add_picture(str(comparison_path), Inches(0.45), Inches(1.15), width=Inches(8.0))
    comp_df = build_comparison_table()
    add_table(slide, 8.65, 1.6, 4.1, 2.2, comp_df, font_size=10)
    add_bullets(
        slide,
        8.7,
        4.2,
        4.0,
        1.8,
        [
            "Wind: hybrid cuts fault error by ~61.5% vs PID.",
            "Sensor: hybrid is slightly better than both baselines on fault-window error.",
            "Comm: hybrid cuts fault error by ~42.9% vs PID.",
        ],
        font_size=15,
    )
    add_footer(slide, "Percentages computed from aggregate fault_error_m values in the final hybrid result package.")

    # Slide 14
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Key Insights", "What the study demonstrates technically")
    add_bullets(
        slide,
        0.75,
        1.4,
        11.3,
        4.9,
        [
            "A layered design works well: keep PID for stabilization and let the supervisor handle context-aware coordination changes.",
            "Supervision helps most when the fault is informational rather than purely dynamic; communication degradation is the best example.",
            "Diagnosis + bounded reconfiguration is preferable to an all-new controller because it preserves transparency and reproducibility.",
            "Estimator quality is the main bottleneck for sensor robustness; this is where future gains are most likely.",
            "A unified hybrid configuration is scientifically stronger than mixing separate result tables from different tuned variants.",
        ],
        font_size=18,
    )
    add_footer(slide, "The final_hybrid package is a single controller/configuration rerun across all scenarios, not a cherry-picked combination of tables.")

    # Slide 15
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Limitations", "What remains to be improved")
    add_bullets(
        slide,
        0.75,
        1.5,
        11.3,
        4.8,
        [
            "No full probabilistic state estimator or bias-state observer is included; the current estimator is a lightweight constant-velocity predictor/corrector.",
            "No formal Lyapunov, ISS, or safety proof is provided for the supervisory switching logic.",
            "The current study uses a 4-UAV square formation only; larger swarms and topology variation remain future work.",
            "Results are simulation-based, so hardware latency, actuation nonlinearities, and aerodynamic coupling are not yet validated experimentally.",
            "The hybrid improves nominal performance too, which changes the paper framing from 'fault-only improvement' to 'overall better coordinated control.'",
        ],
        font_size=17,
    )
    add_footer(slide, "These limitations are acceptable for a simulation paper, but they should be stated explicitly in the manuscript.")

    # Slide 16
    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "Conclusion", "Final takeaways from the study")
    add_bullets(
        slide,
        0.8,
        1.5,
        11.1,
        3.8,
        [
            "A clean standalone simulator was built for diagnosis-driven fault-tolerant UAV swarm control.",
            "The final hybrid controller outperforms PID and the generic supervisor across nominal, wind, sensor, and communication scenarios in the 10-seed package.",
            "The strongest contribution is communication-aware resilience, while sensor robustness benefits from filtered-state handling and remains the main future extension point.",
            "The result package is reproducible, multi-seed validated, and ready to support paper writing and supervisor review.",
        ],
        font_size=18,
    )
    add_text_block(slide, 0.82, 5.85, 11.0, 0.8, f"Final package\nConfig: {CONFIG_PATH}\nResults: {RESULTS_ROOT}", 14, THEME.gray)
    add_footer(slide, "Generated automatically from the paper_sim code and final result files.")

    prs.save(PPTX_PATH)
    return PPTX_PATH


def main() -> None:
    ppt_path = build_presentation()
    print(f"Saved PowerPoint to {ppt_path}")


if __name__ == "__main__":
    main()
